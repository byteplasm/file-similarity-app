import streamlit as st
import os
import shutil
import hashlib
import numpy as np
import pandas as pd
import tempfile
from PyPDF2 import PdfReader
from sentence_transformers import SentenceTransformer
from docx import Document
from pptx import Presentation
from PIL import Image

# ----------------------------
# Functions
# ----------------------------
def get_file_text(file):
    """Extract text for supported file types"""
    try:
        if file.name.endswith(".pdf"):
            reader = PdfReader(file)
            text = "".join(page.extract_text() or "" for page in reader.pages)
            return text
        elif file.name.endswith((".txt", ".csv", ".log")):
            return file.read().decode("utf-8", errors="ignore")
        elif file.name.endswith(".docx"):
            doc = Document(file)
            return "\n".join(p.text for p in doc.paragraphs)
        elif file.name.endswith((".xlsx", ".xls")):
            df = pd.read_excel(file)
            return df.head().to_string()
        elif file.name.endswith(".pptx"):
            prs = Presentation(file)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            return "\n".join(slides_text)
        elif file.name.endswith((".png", ".jpg", ".jpeg", ".gif")):
            return None  # handled as image preview
        else:
            return None
    except:
        return None

def file_hash(file):
    hasher = hashlib.md5()
    try:
        file.seek(0)
        buf = file.read()
        file.seek(0)
        hasher.update(buf)
        return hasher.hexdigest()
    except:
        return None

def similarity(vec1, vec2):
    if vec1 is None or vec2 is None:
        return 0
    return np.dot(vec1, vec2) / (np.linalg.norm(vec1) * np.linalg.norm(vec2))

def group_similar_files(file_list, embeddings, threshold=0.5):
    groups = []
    visited = set()
    for i in range(len(file_list)):
        if file_list[i].name in visited:
            continue
        group = [file_list[i]]
        visited.add(file_list[i].name)
        for j in range(i + 1, len(file_list)):
            if file_list[j].name in visited:
                continue
            score = similarity(embeddings[file_list[i].name], embeddings[file_list[j].name])
            if score >= threshold:
                group.append(file_list[j])
                visited.add(file_list[j].name)
        if len(group) > 1:
            groups.append(group)
    return groups

def auto_rename(temp_folder, file_name):
    base, ext = os.path.splitext(file_name)
    counter = 1
    new_name = f"{base}_duplicate{counter}{ext}"
    new_path = os.path.join(temp_folder, new_name)
    while os.path.exists(new_path):
        counter += 1
        new_name = f"{base}_duplicate{counter}{ext}"
        new_path = os.path.join(temp_folder, new_name)
    return new_path

# ----------------------------
# Load model
# ----------------------------
@st.cache_resource
def load_model():
    return SentenceTransformer('all-MiniLM-L6-v2')

model = load_model()

# ----------------------------
# Streamlit UI
# ----------------------------
st.title("📁 Universal File Duplicate & Similarity Cleaner")

uploaded_files = st.file_uploader(
    "Upload or drag-and-drop files (multiple allowed)",
    type=None,
    accept_multiple_files=True
)

threshold = st.slider("Similarity threshold (text-based files)", 0.0, 1.0, 0.5, 0.01)

cleanup_mode = st.radio(
    "Duplicate handling mode",
    options=[
        "Move duplicates to _duplicates folder",
        "Auto-rename duplicates in place"
    ]
)

if uploaded_files:
    st.write(f"{len(uploaded_files)} files uploaded")
    
    temp_folder = tempfile.mkdtemp()
    file_paths = []
    
    # Save uploaded files to temp folder
    for file in uploaded_files:
        path = os.path.join(temp_folder, file.name)
        with open(path, "wb") as f:
            f.write(file.read())
        file_paths.append(path)
    
    texts = {}
    embeddings = {}
    hashes = {}
    exact_duplicates = []
    
    # Process files
    for file_path in file_paths:
        file = open(file_path, "rb")
        # Exact duplicates via hash
        h = file_hash(file)
        if h:
            if h in hashes:
                exact_duplicates.append((hashes[h], file_path))
            else:
                hashes[h] = file_path
        
        # Text-based embeddings
        file.seek(0)
        text = get_file_text(file)
        if text:
            texts[file_path] = text
            embeddings[file_path] = model.encode(text)
    
    # Show exact duplicates
    st.subheader("🔴 Exact Duplicates")
    if exact_duplicates:
        for a, b in exact_duplicates:
            st.write(f"📄 {os.path.basename(a)}")
            st.write(f"📄 {os.path.basename(b)}")
            st.write("---")
    else:
        st.write("No exact duplicates found.")
    
    # Group similar text files
    if embeddings:
        st.subheader("🟡 Similar Text Files")
        file_objects = [open(p, "rb") for p in file_paths]
        groups = group_similar_files(file_objects, embeddings, threshold)
        
        if not groups:
            st.write("No similar text file groups found above threshold.")
        else:
            decisions = {}
            for idx, group in enumerate(groups):
                st.write(f"### Group {idx+1}")
                keep_file = st.selectbox(
                    "Select file to keep (others handled automatically)",
                    options=[os.path.basename(f.name) for f in group],
                    key=f"keep-{idx}"
                )
                decisions[idx] = {"keep": keep_file, "group": group}
                for jdx, f in enumerate(group):
                    text = texts.get(f.name)
                    if text:
                        with st.expander(f"Preview: {f.name}"):
                            st.text_area("Preview", text, height=150, key=f"preview-{idx}-{jdx}")
                    elif f.name.lower().endswith((".png", ".jpg", ".jpeg", ".gif")):
                        with st.expander(f"Image Preview: {f.name}"):
                            img = Image.open(f)
                            st.image(img, use_column_width=True)
                    else:
                        st.write(f"📄 {f.name} (Preview not available)")

            # Cleanup button
            if st.button("Run Cleanup"):
                archive_folder = os.path.join(temp_folder, "_duplicates")
                os.makedirs(archive_folder, exist_ok=True)
                moved_files = []
                
                # Exact duplicates
                for a, b in exact_duplicates:
                    for f in [a, b]:
                        if cleanup_mode == "Move duplicates to _duplicates folder":
                            dst = os.path.join(archive_folder, os.path.basename(f))
                            shutil.move(f, dst)
                            moved_files.append({"Moved": os.path.basename(f), "To": dst})
                        else:
                            new_path = auto_rename(temp_folder, os.path.basename(f))
                            shutil.move(f, new_path)
                            moved_files.append({"Renamed": os.path.basename(f), "To": new_path})
                
                # Semantic duplicates
                for decision in decisions.values():
                    keep = decision["keep"]
                    for f in decision["group"]:
                        if f.name != keep:
                            if cleanup_mode == "Move duplicates to _duplicates folder":
                                dst = os.path.join(archive_folder, os.path.basename(f.name))
                                shutil.move(f.name, dst)
                                moved_files.append({"Moved": os.path.basename(f.name), "To": dst})
                            else:
                                new_path = auto_rename(temp_folder, os.path.basename(f.name))
                                shutil.move(f.name, new_path)
                                moved_files.append({"Renamed": os.path.basename(f.name), "To": new_path})
                
                if moved_files:
                    st.success(f"Processed {len(moved_files)} files!")
                    st.dataframe(pd.DataFrame(moved_files))
                else:
                    st.info("No files needed processing.")